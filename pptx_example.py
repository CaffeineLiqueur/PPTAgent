"""
python-pptx 操作示例
展示如何使用 python-pptx 库创建和操作 PowerPoint 演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os


def create_basic_presentation():
    """创建基础演示文稿示例"""
    # 创建新的演示文稿
    prs = Presentation()
    
    # 1. 添加标题幻灯片（Title Slide）
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 0 是标题幻灯片布局
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "PPT Agent 演示"
    subtitle1.text = "使用 python-pptx 创建演示文稿"
    
    # 2. 添加标题和内容幻灯片（Title and Content）
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 1 是标题和内容布局
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    title2.text = "项目概述"
    tf = content2.text_frame
    tf.text = "第一点：这是项目的主要特点"
    p = tf.add_paragraph()
    p.text = "第二点：支持多种功能"
    p.level = 1  # 设置缩进级别
    p = tf.add_paragraph()
    p.text = "第三点：易于使用和扩展"
    
    # 3. 添加空白幻灯片并手动添加内容
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # 6 是空白布局
    
    # 添加标题文本框
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(1)
    title_box = slide3.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "自定义内容幻灯片"
    
    # 设置标题样式
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # 添加正文内容
    content_box = slide3.shapes.add_textbox(left, Inches(2), width, Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = "这是一个自定义的文本框，可以设置各种样式和格式。"
    
    # 添加段落
    p = content_frame.add_paragraph()
    p.text = "支持多段落文本"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(51, 51, 51)
    
    # 4. 添加带表格的幻灯片
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "数据表格示例"
    
    # 创建表格
    rows = 4
    cols = 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    table = slide4.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置表格列宽
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(3)
    
    # 填充表头
    table.cell(0, 0).text = "项目"
    table.cell(0, 1).text = "数值"
    table.cell(0, 2).text = "备注"
    
    # 设置表头样式
    for i in range(cols):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
    
    # 填充数据
    data = [
        ["功能A", "85%", "已完成"],
        ["功能B", "60%", "进行中"],
        ["功能C", "30%", "计划中"]
    ]
    
    for i, row_data in enumerate(data, start=1):
        for j, cell_data in enumerate(row_data):
            table.cell(i, j).text = cell_data
            table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 5. 添加带形状的幻灯片
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title5.text_frame.text = "形状和图形示例"
    title5.text_frame.paragraphs[0].font.size = Pt(36)
    title5.text_frame.paragraphs[0].font.bold = True
    
    # 添加矩形
    rect = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0, 102, 204)
    rect.line.color.rgb = RGBColor(0, 51, 102)
    rect.line.width = Pt(2)
    rect.text_frame.text = "矩形"
    rect.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    rect.text_frame.paragraphs[0].font.size = Pt(18)
    rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加圆形
    circle = slide5.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4), Inches(2), Inches(2), Inches(1.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255, 153, 0)
    circle.text_frame.text = "圆形"
    circle.text_frame.paragraphs[0].font.size = Pt(18)
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加箭头
    arrow = slide5.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(7), Inches(2), Inches(2), Inches(1.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0, 153, 76)
    arrow.text_frame.text = "箭头"
    arrow.text_frame.paragraphs[0].font.size = Pt(18)
    arrow.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 6. 添加带图片的幻灯片（需要图片文件存在）
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "图片示例"
    
    # 注意：需要替换为实际的图片路径
    # image_path = "path/to/your/image.jpg"
    # if os.path.exists(image_path):
    #     left = Inches(2)
    #     top = Inches(2)
    #     pic = slide6.shapes.add_picture(image_path, left, top, width=Inches(6))
    
    # 7. 添加带项目符号列表的幻灯片
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "项目列表示例"
    
    content7 = slide7.placeholders[1]
    tf = content7.text_frame
    tf.text = "主要功能"
    
    # 添加多级列表
    items = [
        ("功能模块A", 0),
        ("子功能A1", 1),
        ("子功能A2", 1),
        ("功能模块B", 0),
        ("功能模块C", 0),
    ]
    
    for item_text, level in items:
        p = tf.add_paragraph()
        p.text = item_text
        p.level = level
        p.font.size = Pt(18)
    
    # 8. 设置幻灯片背景
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置背景颜色
    background = slide8.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # 添加内容
    title8 = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title8.text_frame.text = "自定义背景幻灯片"
    title8.text_frame.paragraphs[0].font.size = Pt(36)
    title8.text_frame.paragraphs[0].font.bold = True
    title8.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 9. 添加文本框并设置详细样式
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    
    textbox = slide9.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    
    # 第一段：标题
    p1 = text_frame.paragraphs[0]
    p1.text = "样式丰富的文本"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0, 51, 102)
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(12)
    
    # 第二段：正文
    p2 = text_frame.add_paragraph()
    p2.text = "这是普通正文，可以设置字体大小、颜色等属性。"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(51, 51, 51)
    p2.space_after = Pt(6)
    
    # 第三段：斜体
    p3 = text_frame.add_paragraph()
    p3.text = "这是斜体文本。"
    p3.font.size = Pt(18)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(102, 102, 102)
    p3.space_after = Pt(6)
    
    # 第四段：下划线
    p4 = text_frame.add_paragraph()
    p4.text = "这是带下划线的文本。"
    p4.font.size = Pt(18)
    p4.font.underline = True
    p4.font.color.rgb = RGBColor(0, 102, 204)
    
    # 10. 添加图表幻灯片（饼图示例）
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "图表示例"
    
    # 注意：python-pptx 对图表的支持有限，复杂图表建议使用其他方法
    # 这里展示如何添加占位符说明
    content10 = slide10.placeholders[1]
    tf = content10.text_frame
    tf.text = "注意：python-pptx 对图表的支持有限。"
    p = tf.add_paragraph()
    p.text = "如需复杂图表，可以考虑："
    p = tf.add_paragraph()
    p.text = "1. 使用 matplotlib 生成图片后插入"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. 使用其他库（如 python-pptx-chart）"
    p.level = 1
    
    # 保存演示文稿
    output_path = "example_presentation.pptx"
    prs.save(output_path)
    print(f"演示文稿已保存到: {output_path}")
    
    return prs


def modify_existing_presentation(file_path):
    """修改现有演示文稿的示例"""
    if not os.path.exists(file_path):
        print(f"文件不存在: {file_path}")
        return
    
    prs = Presentation(file_path)
    
    # 获取第一张幻灯片
    if len(prs.slides) > 0:
        first_slide = prs.slides[0]
        
        # 修改标题
        if first_slide.shapes.title:
            first_slide.shapes.title.text = "修改后的标题"
        
        # 添加新幻灯片
        new_slide = prs.slides.add_slide(prs.slide_layouts[1])
        new_slide.shapes.title.text = "新添加的幻灯片"
        new_slide.placeholders[1].text = "这是通过代码添加的内容"
    
    # 保存修改后的文件
    output_path = file_path.replace(".pptx", "_modified.pptx")
    prs.save(output_path)
    print(f"修改后的文件已保存到: {output_path}")


def get_slide_layouts_info(prs):
    """获取所有可用的幻灯片布局信息"""
    print("\n可用的幻灯片布局:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")


if __name__ == "__main__":
    # 创建新演示文稿
    print("正在创建演示文稿...")
    prs = create_basic_presentation()
    
    # 显示可用的布局
    get_slide_layouts_info(prs)
    
    print("\n示例完成！")
    print("运行此脚本将生成一个包含多种元素类型的示例 PPT 文件。")

